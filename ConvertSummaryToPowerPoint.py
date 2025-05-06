from pptx import Presentation
from pptx.util import Pt
from tkinter import filedialog, Tk
import os
import json
import openai


openai.api_key = os.getenv('OPENAI_API_KEY')

def set_font_size(shape, size=16):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    set_font_size(slide.placeholders[1], 16)

def add_two_column_slide(prs, content_dict, title="Slide", left_keys=3, format_right=None):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = title
    left_text = "\n".join(f"{k}: {v}" for k, v in list(content_dict.items())[:left_keys])
    if format_right:
        right_text = format_right(content_dict)
    else:
        right_text = "\n".join(f"{k}: {v}" for k, v in list(content_dict.items())[left_keys:])
    slide.placeholders[1].text = left_text or "Not found."
    slide.placeholders[2].text = right_text or "Not found."
    set_font_size(slide.placeholders[1], 16)
    set_font_size(slide.placeholders[2], 16)

def generate_slide_content_from_summary(summary_text):
    prompt = f"""
You are a neuroscience research assistant preparing a PowerPoint presentation based on the summary of a scientific paper.

The summary below was generated from a neuroscience article. Reformat it into structured content for a 4-slide presentation.

Use the following JSON format:
{{
  "title_slide": "Summary of: [Paper Title]",
  "metadata": {{
    "Paper Name": "...",
    "Authors": "...",
    "Year": "...",
    "Journal": "...",
    "Impact Factor": "...",
    "Keywords": "..."
  }},
  "experimental_design": {{
    "Tasks": "...",
    "Modalities": "...",
    "Regions": "...",
    "Sample Size": "...",
    "Trial Structure": "...",
    "Statistical Comparisons": "..."
  }},
  "results_and_conclusion": {{
    "Findings": "...",
    "Conclusions": "...",
    "Future Work": "..."
  }}
}}

ONLY return a valid JSON object — no markdown, no commentary.

SUMMARY:
\"\"\"
{summary_text}
\"\"\"
"""

    response = openai.ChatCompletion.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a structured neuroscience research assistant preparing slide content."},
            {"role": "user", "content": prompt}
        ],
        temperature=0,
        max_tokens=2048
    )

    content = response.choices[0].message.content.strip()

    # Remove Markdown code block (```json ... ```)
    if content.startswith("```json"):
        content = content.replace("```json", "").replace("```", "").strip()
    elif content.startswith("```"):
        content = content.replace("```", "").strip()

    try:
        return json.loads(content)
    except json.JSONDecodeError as e:
        print("JSON decoding error after cleaning:", e)
        return None

def create_ppt_from_summaries(summary_dir):
    prs = Presentation()
    add_title_slide(prs, "EEG Summaries Presentation", "A_Cerullo")

    # Process Proposed_Hypothesis_And_Protocol.txt
    hypothesis_path = os.path.join(summary_dir, "Proposed_Hypothesis_And_Protocol.txt")
    if os.path.exists(hypothesis_path):
        print("\nProcessing: Proposed_Hypothesis_And_Protocol.txt")
        with open(hypothesis_path, "r", encoding="utf-8") as f:
            hypothesis_text = f.read()

        # Try splitting based on section headers or fallback to truncation
        sections = hypothesis_text.split("Experimental Design:")
        if len(sections) == 2:
            intro, design = sections
        else:
            intro = hypothesis_text[:1000]
            design = hypothesis_text[1000:]

        # Slide 1: Hypothesis and Background
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])
        slide1.shapes.title.text = "Proposed Hypothesis & Justification"
        slide1.placeholders[1].text = intro.strip()
        set_font_size(slide1.placeholders[1], 16)

        # Slide 2: Experimental Design
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Experimental Design"
        slide2.placeholders[1].text = "Experimental Design: " + design.strip()
        set_font_size(slide2.placeholders[1], 16)
    else:
        print("Proposed_Hypothesis_And_Protocol.txt not found.")

    # Process AllSummaries_MetaSummary.txt
    meta_summary_path = os.path.join(summary_dir, "AllSummaries_MetaSummary.txt")
    if os.path.exists(meta_summary_path):
        print("\nProcessing: AllSummaries_MetaSummary.txt")
        with open(meta_summary_path, "r", encoding="utf-8") as f:
            meta_summary_text = f.read()

        slide_data = generate_slide_content_from_summary(meta_summary_text)

        if slide_data:
            add_title_slide(prs, slide_data.get("title_slide", "Meta Summary"), "")
            add_two_column_slide(prs, slide_data.get("metadata", {}), title="Meta Summary: Metadata")
            add_two_column_slide(prs, slide_data.get("experimental_design", {}), title="Meta Summary: Experimental Design")
            add_two_column_slide(prs, slide_data.get("results_and_conclusion", {}), title="Meta Summary: Results and Conclusion", left_keys=2, format_right=lambda d: f"Future Work: {d.get('Future Work', 'N/A')}")
        else:
            print("Could not generate slide content from AllSummaries_MetaSummary.txt.")
    else:
        print("AllSummaries_MetaSummary.txt not found.")


    add_title_slide(prs, "Individual Paper Summaries", "")

    summary_files = sorted(f for f in os.listdir(summary_dir) if f.endswith("_summary.txt"))

    for file in summary_files:
        path = os.path.join(summary_dir, file)
        with open(path, 'r', encoding='utf-8') as f:
            summary_text = f.read()

        print(f"Processing: {file}")
        slide_data = generate_slide_content_from_summary(summary_text)

        if not slide_data:
            print(f"Skipping {file} due to parse failure.")
            continue

        add_title_slide(prs, slide_data.get("title_slide", "Summary"), "")
        add_two_column_slide(prs, slide_data.get("metadata", {}), title="Metadata")
        add_two_column_slide(prs, slide_data.get("experimental_design", {}), title="Experimental Design")
        add_two_column_slide(prs, slide_data.get("results_and_conclusion", {}), title="Results and Conclusion", left_keys=2, format_right=lambda d: f"Future Work: {d.get('Future Work', 'N/A')}")

    
    


    output_path = os.path.join(summary_dir, "EEG_Summaries_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

def main():
    root = Tk()
    root.withdraw()
    selected_dir = filedialog.askdirectory(title="Select Folder Containing Summary Files")
    if not selected_dir:
        print("No folder selected.")
        return

    create_ppt_from_summaries(selected_dir)

if __name__ == "__main__":
    main()
